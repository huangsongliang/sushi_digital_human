"""Office 文档解析器
支持 Word (.docx)、PowerPoint (.pptx) 文档解析
"""

from pathlib import Path
from typing import Any, Dict, List, Optional

from backend.utils.logger import get_logger

logger = get_logger(__name__)


class WordParser:
    """Word 文档解析器"""

    def __init__(self):
        self._engine = None

    def _init_engine(self):
        """延迟初始化"""
        if self._engine is None:
            try:
                import docx

                self._engine = docx
                logger.info("Word 解析引擎初始化成功")
            except ImportError:
                logger.warning("python-docx 未安装，请安装: pip install python-docx")
                self._engine = False

    def extract_text(self, file_path: str) -> str:
        """提取 Word 文档文本"""
        self._init_engine()

        if not self._engine:
            raise RuntimeError("Word 解析引擎未初始化，请安装 python-docx")

        try:
            doc = self._engine.Document(file_path)
            paragraphs = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)

            full_text = "\n\n".join(paragraphs)
            logger.info(f"Word 文档文本提取完成: {len(full_text)} 字符")
            return full_text

        except Exception as e:
            logger.error(f"Word 文档解析失败: {str(e)}")
            raise

    def extract_tables(self, file_path: str) -> List[Dict[str, Any]]:
        """提取 Word 文档中的表格"""
        self._init_engine()

        if not self._engine:
            return []

        try:
            doc = self._engine.Document(file_path)
            tables = []

            for table_idx, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)

                tables.append(
                    {
                        "table_index": table_idx,
                        "rows": rows,
                        "row_count": len(rows),
                        "col_count": len(rows[0]) if rows else 0,
                    }
                )

            logger.info(f"Word 表格提取完成: {len(tables)} 个表格")
            return tables

        except Exception as e:
            logger.error(f"Word 表格提取失败: {str(e)}")
            return []

    def extract_images(self, file_path: str, output_dir: Optional[str] = None) -> List[str]:
        """提取 Word 文档中的图片"""
        self._init_engine()

        if not self._engine:
            return []

        try:
            doc = self._engine.Document(file_path)

            if output_dir is None:
                output_dir = Path(file_path).parent / f"{Path(file_path).stem}_images"
            else:
                output_dir = Path(output_dir)

            output_dir.mkdir(parents=True, exist_ok=True)

            image_paths = []
            image_counter = 0

            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image = rel.target_part.blob
                    image_ext = rel.target_ref.split(".")[-1]
                    image_name = f"word_image_{image_counter}.{image_ext}"
                    image_path = output_dir / image_name

                    with open(image_path, "wb") as f:
                        f.write(image)

                    image_paths.append(str(image_path))
                    image_counter += 1

            logger.info(f"Word 图片提取完成: {len(image_paths)} 张图片")
            return image_paths

        except Exception as e:
            logger.error(f"Word 图片提取失败: {str(e)}")
            return []

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """提取 Word 文档元数据"""
        self._init_engine()

        if not self._engine:
            return {}

        try:
            doc = self._engine.Document(file_path)
            core_props = doc.core_properties

            return {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "paragraphs": len([p for p in doc.paragraphs if p.text.strip()]),
                "tables": len(doc.tables),
            }

        except Exception as e:
            logger.error(f"Word 元数据提取失败: {str(e)}")
            return {}

    def process(self, file_path: str, extract_tables: bool = True, extract_images: bool = False) -> Dict[str, Any]:
        """完整处理 Word 文档"""
        result = {
            "file_path": file_path,
            "file_name": Path(file_path).name,
            "text": self.extract_text(file_path),
            "metadata": self.extract_metadata(file_path),
            "tables": [],
            "images": [],
        }

        if extract_tables:
            result["tables"] = self.extract_tables(file_path)

        if extract_images:
            result["images"] = self.extract_images(file_path)

        return result


class PowerPointParser:
    """PowerPoint 文档解析器"""

    def __init__(self):
        self._engine = None

    def _init_engine(self):
        """延迟初始化"""
        if self._engine is None:
            try:
                import pptx

                self._engine = pptx
                logger.info("PowerPoint 解析引擎初始化成功")
            except ImportError:
                logger.warning("python-pptx 未安装，请安装: pip install python-pptx")
                self._engine = False

    def extract_text(self, file_path: str) -> str:
        """提取 PowerPoint 文本"""
        self._init_engine()

        if not self._engine:
            raise RuntimeError("PowerPoint 解析引擎未初始化，请安装 python-pptx")

        try:
            prs = self._engine.Presentation(file_path)
            slides_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"[幻灯片 {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if len(slide_texts) > 1:
                    slides_text.append("\n".join(slide_texts))

            full_text = "\n\n".join(slides_text)
            logger.info(f"PowerPoint 文本提取完成: {len(full_text)} 字符")
            return full_text

        except Exception as e:
            logger.error(f"PowerPoint 文档解析失败: {str(e)}")
            raise

    def extract_notes(self, file_path: str) -> List[Dict[str, Any]]:
        """提取幻灯片备注"""
        self._init_engine()

        if not self._engine:
            return []

        try:
            prs = self._engine.Presentation(file_path)
            notes_list = []

            for slide_num, slide in enumerate(prs.slides, 1):
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_list.append(
                            {
                                "slide": slide_num,
                                "notes": notes_text,
                            }
                        )

            return notes_list

        except Exception as e:
            logger.error(f"PowerPoint 备注提取失败: {str(e)}")
            return []

    def extract_images(self, file_path: str, output_dir: Optional[str] = None) -> List[str]:
        """提取 PowerPoint 图片"""
        self._init_engine()

        if not self._engine:
            return []

        try:
            prs = self._engine.Presentation(file_path)

            if output_dir is None:
                output_dir = Path(file_path).parent / f"{Path(file_path).stem}_images"
            else:
                output_dir = Path(output_dir)

            output_dir.mkdir(parents=True, exist_ok=True)

            image_paths = []
            image_counter = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_ext = image.ext
                        image_name = f"ppt_slide{slide_num}_image{image_counter}.{image_ext}"
                        image_path = output_dir / image_name

                        with open(image_path, "wb") as f:
                            f.write(image.blob)

                        image_paths.append(str(image_path))
                        image_counter += 1

            logger.info(f"PowerPoint 图片提取完成: {len(image_paths)} 张图片")
            return image_paths

        except Exception as e:
            logger.error(f"PowerPoint 图片提取失败: {str(e)}")
            return []

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """提取 PowerPoint 元数据"""
        self._init_engine()

        if not self._engine:
            return {}

        try:
            prs = self._engine.Presentation(file_path)

            return {
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "subject": prs.core_properties.subject or "",
                "slide_count": len(prs.slides),
                "created": str(prs.core_properties.created) if prs.core_properties.created else "",
                "modified": str(prs.core_properties.modified) if prs.core_properties.modified else "",
            }

        except Exception as e:
            logger.error(f"PowerPoint 元数据提取失败: {str(e)}")
            return {}

    def process(self, file_path: str, extract_notes: bool = True, extract_images: bool = False) -> Dict[str, Any]:
        """完整处理 PowerPoint 文档"""
        result = {
            "file_path": file_path,
            "file_name": Path(file_path).name,
            "text": self.extract_text(file_path),
            "metadata": self.extract_metadata(file_path),
            "notes": [],
            "images": [],
        }

        if extract_notes:
            result["notes"] = self.extract_notes(file_path)

        if extract_images:
            result["images"] = self.extract_images(file_path)

        return result


def parse_word_document(file_path: str) -> Dict[str, Any]:
    """便捷函数：解析 Word 文档"""
    parser = WordParser()
    return parser.process(file_path, extract_tables=True, extract_images=False)


def parse_powerpoint_document(file_path: str) -> Dict[str, Any]:
    """便捷函数：解析 PowerPoint 文档"""
    parser = PowerPointParser()
    return parser.process(file_path, extract_notes=True, extract_images=False)
